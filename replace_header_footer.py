import os
from pathlib import Path
from docx import Document
from docx.shared import Inches
import subprocess
import platform
import shutil

# Configuration des chemins
SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
INPUT_FOLDER = os.path.join(SCRIPT_DIR, "input_docs")
OUTPUT_FOLDER = os.path.join(SCRIPT_DIR, "output_docs")
PDF_OUTPUT_FOLDER = os.path.join(SCRIPT_DIR, "pdf_output")
LOGO_PATH = os.path.join(SCRIPT_DIR, "logo.png")
FOOTER_PATH = os.path.join(SCRIPT_DIR, "footer.txt")

# Extensions supportées
SUPPORTED_EXTENSIONS = {
    '.docx': 'Document Word',
    '.pptx': 'PowerPoint',
    '.xlsx': 'Excel'
}

def setup_folders():
    """Crée les dossiers nécessaires s'ils n'existent pas"""
    for folder in [INPUT_FOLDER, OUTPUT_FOLDER, PDF_OUTPUT_FOLDER]:
        os.makedirs(folder, exist_ok=True)
        print(f"✓ Dossier créé/vérifié : {folder}")

def get_footer_text():
    """Lit le texte du pied de page depuis le fichier footer.txt"""
    try:
        with open(FOOTER_PATH, 'r', encoding='utf-8') as f:
            return f.read().strip()
    except FileNotFoundError:
        print(f"⚠️  Le fichier footer.txt n'a pas été trouvé à l'emplacement : {FOOTER_PATH}")
        print("➡️  Création du fichier footer.txt avec le texte par défaut")
        default_footer = "© 2024 Management Solution | Mentions légales"
        with open(FOOTER_PATH, 'w', encoding='utf-8') as f:
            f.write(default_footer)
        return default_footer

def process_document(doc_path, output_path, footer_text):
    """Traite un document Word"""
    doc = Document(doc_path)
    
    for section in doc.sections:
        # Modification de l'en-tête
        header = section.header
        for paragraph in header.paragraphs:
            paragraph.clear()
        if not header.paragraphs:
            header.add_paragraph()
        run = header.paragraphs[0].add_run()
        run.add_picture(LOGO_PATH, width=Inches(1.5))
        
        # Modification du pied de page
        footer = section.footer
        for paragraph in footer.paragraphs:
            paragraph.clear()
        if not footer.paragraphs:
            footer.add_paragraph()
        
        # Ajout d'espaces pour positionner le pied de page plus bas
        for _ in range(2):  # Ajout de paragraphes vides pour pousser le contenu vers le bas
            footer.add_paragraph()
        
        # Ajout du texte du pied de page dans le dernier paragraphe
        last_paragraph = footer.paragraphs[-1]
        last_paragraph.text = footer_text
    
    doc.save(output_path)

def process_file(input_path, output_path, file_type, footer_text):
    """Traite un fichier selon son type"""
    try:
        if file_type == '.docx':
            process_document(input_path, output_path, footer_text)
        elif file_type == '.xlsx':
            # Pour les fichiers Excel, on copie simplement le fichier pour l'instant
            # car la bibliothèque python-docx ne gère pas les fichiers Excel
            shutil.copy2(input_path, output_path)
        return True
    except Exception as e:
        print(f"❌ Erreur lors du traitement de {input_path}: {str(e)}")
        return False

def replace_variables_in_document(doc_path, output_path, client_data):
    """Remplace les variables dans un document Word par les données du client"""
    doc = Document(doc_path)
    
    # Parcourir tous les paragraphes du document
    for paragraph in doc.paragraphs:
        for key, value in client_data.items():
            if f"«{key}»" in paragraph.text:
                paragraph.text = paragraph.text.replace(f"«{key}»", value)
    
    # Parcourir toutes les tables
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                for paragraph in cell.paragraphs:
                    for key, value in client_data.items():
                        if f"«{key}»" in paragraph.text:
                            paragraph.text = paragraph.text.replace(f"«{key}»", value)
    
    # Parcourir les en-têtes et pieds de page
    for section in doc.sections:
        # En-têtes
        header = section.header
        for paragraph in header.paragraphs:
            for key, value in client_data.items():
                if f"«{key}»" in paragraph.text:
                    paragraph.text = paragraph.text.replace(f"«{key}»", value)
        
        # Pieds de page
        footer = section.footer
        for paragraph in footer.paragraphs:
            for key, value in client_data.items():
                if f"«{key}»" in paragraph.text:
                    paragraph.text = paragraph.text.replace(f"«{key}»", value)
    
    # Sauvegarder le document modifié
    doc.save(output_path)
    return True

def process_client_template(input_path, output_path, file_type, client_data, footer_text=None, logo_path=None):
    """Traite un template client selon son type"""
    try:
        if file_type == '.docx':
            # Créer une copie temporaire du fichier d'entrée
            temp_file = output_path + ".tmp.docx"
            shutil.copy2(input_path, temp_file)
            
            # Remplacer les variables
            replace_variables_in_document(temp_file, output_path, client_data)
            
            # Vérifier si le fichier a été modifié
            if not os.path.exists(output_path):
                raise Exception("Le fichier n'a pas été créé")
            
            # Ajouter le logo et le pied de page si demandé
            if footer_text or (logo_path and os.path.exists(logo_path)):
                doc = Document(output_path)
                
                for section in doc.sections:
                    # Modifier l'en-tête si un logo est fourni et existe
                    if logo_path and os.path.exists(logo_path):
                        header = section.header
                        for paragraph in header.paragraphs:
                            paragraph.clear()
                        if not header.paragraphs:
                            header.add_paragraph()
                        run = header.paragraphs[0].add_run()
                        try:
                            run.add_picture(logo_path, width=Inches(1.5))
                        except Exception as e:
                            print(f"⚠️ Erreur lors de l'ajout du logo : {str(e)}")
                    
                    # Modifier le pied de page si demandé
                    if footer_text:
                        footer = section.footer
                        for paragraph in footer.paragraphs:
                            paragraph.clear()
                        if not footer.paragraphs:
                            footer.add_paragraph()
                        
                        # Ajout d'espaces pour positionner le pied de page plus bas
                        for _ in range(2):
                            footer.add_paragraph()
                        
                        # Ajout du texte du pied de page
                        last_paragraph = footer.paragraphs[-1]
                        last_paragraph.text = footer_text
                
                doc.save(output_path)
            
            # Nettoyer le fichier temporaire
            if os.path.exists(temp_file):
                os.remove(temp_file)
                
        elif file_type == '.xlsx':
            # Pour Excel, on va créer un fichier temporaire mais avec la bonne extension
            temp_file = output_path + ".tmp.xlsx"
            shutil.copy2(input_path, temp_file)
            
            # Remplacer les variables dans le fichier Excel
            try:
                import openpyxl
                from openpyxl.drawing.image import Image
                
                wb = openpyxl.load_workbook(temp_file)
                
                # Parcourir toutes les feuilles
                for sheet_name in wb.sheetnames:
                    ws = wb[sheet_name]
                    
                    # Ajouter le logo à la première feuille seulement
                    if sheet_name == wb.sheetnames[0] and logo_path and os.path.exists(logo_path):
                        try:
                            # Ajouter le logo dans la cellule A1
                            img = Image(logo_path)
                            # Redimensionner l'image
                            img.width = 150
                            img.height = 75
                            ws.add_image(img, "A1")
                        except Exception as e:
                            print(f"⚠️ Erreur lors de l'ajout du logo dans Excel : {str(e)}")
                    
                    # Parcourir toutes les cellules pour remplacer les variables
                    for row in ws.rows:
                        for cell in row:
                            if isinstance(cell.value, str):
                                for key, value in client_data.items():
                                    if f"«{key}»" in cell.value:
                                        cell.value = cell.value.replace(f"«{key}»", str(value))
                
                # Ajouter le pied de page si fourni
                if footer_text:
                    for ws in wb.worksheets:
                        try:
                            # Excel n'a pas de pied de page facilement accessible via l'API
                            # On ajoute un texte dans les dernières lignes de la première feuille
                            last_row = ws.max_row + 2
                            ws.cell(row=last_row, column=1).value = footer_text
                        except Exception as e:
                            print(f"⚠️ Erreur lors de l'ajout du pied de page dans Excel : {str(e)}")
                
                wb.save(output_path)
            except Exception as e:
                print(f"⚠️ Erreur lors du traitement du fichier Excel {input_path}: {str(e)}")
                # En cas d'erreur, on copie simplement le fichier original
                shutil.copy2(input_path, output_path)
            
            # Nettoyer le fichier temporaire
            if os.path.exists(temp_file):
                os.remove(temp_file)
        
        elif file_type == '.pptx':
            # Pour PowerPoint, on crée un fichier temporaire avec la bonne extension
            temp_file = output_path + ".tmp.pptx"
            shutil.copy2(input_path, temp_file)
            
            try:
                from pptx import Presentation
                from pptx.util import Inches
                
                # Charger la présentation
                prs = Presentation(temp_file)
                
                # Ajouter le logo à toutes les diapositives
                if logo_path and os.path.exists(logo_path):
                    try:
                        for slide in prs.slides:
                            # Ajouter le logo en haut à gauche
                            left = Inches(0.5)
                            top = Inches(0.5)
                            width = Inches(1.5)
                            slide.shapes.add_picture(logo_path, left, top, width=width)
                    except Exception as e:
                        print(f"⚠️ Erreur lors de l'ajout du logo dans PowerPoint : {str(e)}")
                
                # Parcourir toutes les diapositives
                for slide in prs.slides:
                    # Parcourir tous les shapes (zones de texte, etc.)
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            # Remplacer les variables dans le texte
                            text = shape.text
                            for key, value in client_data.items():
                                if f"«{key}»" in text:
                                    text = text.replace(f"«{key}»", str(value))
                            shape.text = text
                    
                    # Ajouter le pied de page sur chaque diapositive
                    if footer_text:
                        try:
                            # Dimensions de la diapositive
                            slide_width = prs.slide_width
                            slide_height = prs.slide_height
                            
                            # Ajouter une zone de texte pour le pied de page
                            left = Inches(0.5)
                            top = slide_height - Inches(1)
                            width = slide_width - Inches(1)
                            height = Inches(0.8)
                            
                            textbox = slide.shapes.add_textbox(left, top, width, height)
                            textbox.text = footer_text
                            textbox.text_frame.paragraphs[0].font.size = Inches(0.1)
                        except Exception as e:
                            print(f"⚠️ Erreur lors de l'ajout du pied de page dans PowerPoint : {str(e)}")
                
                # Sauvegarder la présentation
                prs.save(output_path)
            except Exception as e:
                print(f"⚠️ Erreur lors du traitement du fichier PowerPoint {input_path}: {str(e)}")
                # En cas d'erreur, on garde la copie simple
                shutil.copy2(input_path, output_path)
            
            # Nettoyer le fichier temporaire
            if os.path.exists(temp_file):
                os.remove(temp_file)
        
        else:
            # Pour les autres types de fichiers, faire une copie simple
            shutil.copy2(input_path, output_path)
            
        return True
    except Exception as e:
        print(f"❌ Erreur lors du traitement de {input_path}: {str(e)}")
        return False

# Création des dossiers nécessaires
for folder in [INPUT_FOLDER, OUTPUT_FOLDER, PDF_OUTPUT_FOLDER]:
    os.makedirs(folder, exist_ok=True)

# ... reste du code ...